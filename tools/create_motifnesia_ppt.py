from pathlib import Path

import fitz
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PDF = Path(r"D:\Telyu\Semester 4\Proyek SI\coba.pdf")
OUT_DIR = ROOT / "generated"
ASSET_DIR = OUT_DIR / "ppt_assets"
PPTX = OUT_DIR / "MotifNesia_Outline_Presentasi.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(37, 32, 29)
MUTED = RGBColor(111, 94, 85)
CREAM = RGBColor(250, 246, 238)
PAPER = RGBColor(255, 252, 246)
TERRACOTTA = RGBColor(166, 83, 51)
INDIGO = RGBColor(33, 61, 102)
GOLD = RGBColor(202, 154, 77)
GREEN = RGBColor(69, 116, 94)
WHITE = RGBColor(255, 255, 255)


def ensure_dirs():
    OUT_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(exist_ok=True)


def autocrop(path: Path) -> None:
    im = Image.open(path).convert("RGB")
    bg = Image.new("RGB", im.size, (255, 255, 255))
    diff = ImageChops.difference(im, bg)
    bbox = diff.getbbox()
    if bbox:
        pad = 24
        left = max(bbox[0] - pad, 0)
        top = max(bbox[1] - pad, 0)
        right = min(bbox[2] + pad, im.width)
        bottom = min(bbox[3] + pad, im.height)
        im.crop((left, top, right, bottom)).save(path)


def render_page(doc, page_no: int, name: str) -> Path:
    path = ASSET_DIR / f"{name}.png"
    page = doc[page_no - 1]
    pix = page.get_pixmap(matrix=fitz.Matrix(2.2, 2.2), alpha=False)
    pix.save(path)
    autocrop(path)
    return path


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    return shape


def set_text(tf, text, size=22, color=INK, bold=False, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align


def add_title(slide, title, kicker=None):
    if kicker:
        k = slide.shapes.add_textbox(Inches(0.68), Inches(0.35), Inches(5.5), Inches(0.28))
        set_text(k.text_frame, kicker.upper(), 8.5, TERRACOTTA, True)
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.62), Inches(8.3), Inches(0.68))
    set_text(box.text_frame, title, 26, INK, True)
    add_rect(slide, Inches(0.65), Inches(1.28), Inches(1.0), Inches(0.05), GOLD)


def add_footer(slide, num):
    box = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.45), Inches(0.22))
    set_text(box.text_frame, f"{num:02d}", 8, MUTED, True, PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=17, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_image(slide, path, x, y, w, h, caption=None):
    add_rect(slide, x - Inches(0.04), y - Inches(0.04), w + Inches(0.08), h + Inches(0.08), WHITE, RGBColor(235, 225, 214))
    im = Image.open(path)
    img_ratio = im.width / im.height
    box_ratio = w / h
    if img_ratio > box_ratio:
        pic_w = w
        pic_h = w / img_ratio
        px = x
        py = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        px = x + (w - pic_w) / 2
        py = y
    slide.shapes.add_picture(str(path), px, py, width=pic_w, height=pic_h)
    if caption:
        cap = slide.shapes.add_textbox(x, y + h + Inches(0.08), w, Inches(0.28))
        set_text(cap.text_frame, caption, 8.5, MUTED)


def blank_slide(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_rect(slide, Inches(0), Inches(0), Inches(0.16), SLIDE_H, TERRACOTTA)
    add_footer(slide, num)
    return slide


def build_deck():
    ensure_dirs()
    doc = fitz.open(PDF)
    assets = {
        "as_is": render_page(doc, 23, "as_is_gudang_supplier"),
        "to_be": render_page(doc, 31, "to_be_verifikasi_pembayaran"),
        "persona": render_page(doc, 38, "journey_persona_udin"),
        "journey": render_page(doc, 39, "journey_persona_aulia"),
        "architecture": render_page(doc, 53, "arsitektur_sistem"),
        "usecase": render_page(doc, 57, "diagram_use_case"),
        "catalog": render_page(doc, 126, "ui_beranda_katalog"),
        "report": render_page(doc, 139, "ui_laporan_penjualan"),
        "stock": render_page(doc, 146, "ui_gudang_stok"),
        "supplier": render_page(doc, 152, "ui_dashboard_supplier"),
    }

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1
    s = blank_slide(prs, 1)
    add_rect(s, Inches(0.65), Inches(0.75), Inches(5.0), Inches(0.08), GOLD)
    t = s.shapes.add_textbox(Inches(0.65), Inches(1.02), Inches(6.7), Inches(1.45))
    set_text(t.text_frame, "Pengembangan Website Batik MotifNesia", 36, INK, True)
    st = s.shapes.add_textbox(Inches(0.68), Inches(2.55), Inches(6.2), Inches(0.9))
    set_text(st.text_frame, "Berbasis Web sebagai Media Penjualan Online untuk Toko Viera Garut", 18, MUTED)
    add_bullets(
        s,
        [
            "Proyek Sistem Informasi",
            "D3 Sistem Informasi - Fakultas Ilmu Terapan",
            "Universitas Telkom, 2026",
        ],
        Inches(0.7),
        Inches(4.65),
        Inches(5.6),
        Inches(1.1),
        15,
        INK,
    )
    add_image(s, assets["catalog"], Inches(7.3), Inches(0.75), Inches(5.15), Inches(5.8), "Cuplikan antarmuka katalog MotifNesia")

    slides = [
        (
            "Latar Belakang",
            [
                "Batik memiliki nilai budaya dan ekonomi tinggi, sekaligus menjadi bagian penting dari UMKM Indonesia.",
                "Toko Viera Garut masih mengandalkan toko fisik dan promosi dari mulut ke mulut.",
                "Kondisi ini membatasi jangkauan pemasaran, informasi produk, dan daya saing di pasar digital.",
            ],
            None,
            None,
        ),
        (
            "Permasalahan Utama",
            [
                "Belum tersedia media penjualan online yang terintegrasi.",
                "Informasi produk belum tersaji lengkap dan mudah diakses pelanggan.",
                "Transaksi, stok, retur, dan laporan masih banyak dilakukan secara manual.",
                "Pemesanan via WhatsApp belum mendukung pelacakan pesanan secara mandiri.",
            ],
            assets["as_is"],
            "Contoh proses manual: gudang dan pemesanan ke supplier",
        ),
        (
            "Tujuan Proyek",
            [
                "Merancang website MotifNesia sesuai kebutuhan Toko Viera Garut.",
                "Membangun e-commerce untuk memperluas pemasaran produk batik.",
                "Menyediakan katalog produk yang jelas, informatif, dan mudah dicari.",
                "Mengimplementasikan transaksi online agar pemesanan lebih efektif dan efisien.",
            ],
            None,
            None,
        ),
        (
            "Kondisi Sistem Saat Ini",
            [
                "Pelanggan datang langsung ke toko atau menghubungi melalui WhatsApp.",
                "Pengecekan stok dilakukan oleh petugas berdasarkan catatan atau kondisi fisik barang.",
                "Laporan penjualan dan stock opname membutuhkan rekap manual.",
                "Retur dan pengadaan supplier belum terdokumentasi dalam satu sistem digital.",
            ],
            assets["as_is"],
            "Alur bisnis as-is yang masih bergantung pada proses manual",
        ),
        (
            "Kesenjangan Proses Bisnis",
            [
                "Informasi stok lambat dan berisiko tidak akurat.",
                "Pesanan yang belum dibayar dapat menahan stok terlalu lama.",
                "Pemilik toko sulit memperoleh laporan real-time untuk pengambilan keputusan.",
                "Pelanggan membutuhkan transparansi produk, pembayaran, dan pengiriman.",
            ],
            None,
            None,
        ),
        (
            "Solusi To-Be MotifNesia",
            [
                "Pemesanan digital terintegrasi dari katalog hingga pengiriman.",
                "Inventaris real-time dengan notifikasi stok minimum.",
                "Verifikasi pembayaran dan pembatalan otomatis untuk pesanan kedaluwarsa.",
                "Retur digital, live chat, review produk, dan membership pelanggan.",
            ],
            assets["to_be"],
            "Proses verifikasi pembayaran dan pengelolaan pesanan expired",
        ),
        (
            "Analisis Pengguna",
            [
                "Customer membutuhkan katalog jelas, transaksi praktis, dan pelacakan pesanan.",
                "Admin membutuhkan kendali produk, pesanan, laporan, retur, stok, dan promosi.",
                "Supplier membutuhkan akses untuk memantau pengadaan stok.",
                "Persona utama: pekerja profesional dan mahasiswa dengan kebutuhan belanja yang berbeda.",
            ],
            assets["persona"],
            "User journey persona pekerja profesional",
        ),
        (
            "Kebutuhan Fungsional",
            [
                "Registrasi, login, lupa password, profil, dan notifikasi.",
                "Homepage, katalog produk, filter, detail produk, favorit, dan keranjang.",
                "Checkout, metode pembayaran, konfirmasi pembayaran, riwayat pembelian, dan tracking.",
                "Ulasan produk, promo, voucher, dan membership.",
            ],
            assets["journey"],
            "Kebutuhan fitur dirancang dari perjalanan pengguna",
        ),
        (
            "Kebutuhan Operasional Admin",
            [
                "Kelola produk, kategori, harga, stok, gambar, dan konten statis.",
                "Verifikasi pembayaran, pembaruan status pesanan, dan pelacakan pengiriman.",
                "Laporan penjualan, manajemen membership, voucher, broadcast, dan review.",
                "Manajemen retur, supplier, pengadaan stok, gudang, dan stock opname.",
            ],
            assets["report"],
            "Dashboard laporan penjualan untuk admin",
        ),
        (
            "Kebutuhan Non-Fungsional",
            [
                "Data produk, transaksi, stok, dan pesanan harus akurat serta mudah dipantau.",
                "Keamanan akun, data pribadi, dan transaksi menjadi prioritas utama.",
                "Antarmuka perlu responsif untuk desktop maupun perangkat mobile.",
                "Sistem harus stabil agar proses belanja dan operasional toko tidak terhambat.",
            ],
            None,
            None,
        ),
        (
            "Arsitektur Sistem",
            [
                "Frontend web dan mobile mengakses sistem melalui lapisan API.",
                "Backend service mengelola autentikasi, produk, keranjang, checkout, pembayaran, pesanan, notifikasi, dan laporan.",
                "Database menyimpan data pengguna, produk, transaksi, stok, membership, dan supplier.",
                "Arsitektur ini memisahkan tampilan, logika bisnis, dan penyimpanan data.",
            ],
            assets["architecture"],
            "Arsitektur sistem MotifNesia",
        ),
        (
            "Perancangan Sistem",
            [
                "Pemodelan mencakup hierarki pengguna, use case, activity diagram, ERD, dan class diagram.",
                "Aktor utama sistem adalah Guest, Customer, Admin, dan Supplier.",
                "Use case menggambarkan hak akses serta aktivitas tiap aktor dalam sistem.",
                "Perancangan menjadi dasar implementasi fitur pada web MotifNesia.",
            ],
            assets["usecase"],
            "Diagram use case sistem MotifNesia",
        ),
        (
            "Implementasi Antarmuka",
            [
                "Customer dapat mengakses katalog, detail produk, keranjang, checkout, pembayaran, membership, retur, dan live chat.",
                "Admin mengelola produk, pesanan, laporan, membership, review, retur, gudang, dan supplier.",
                "Supplier dapat memantau dan memperbarui status pengadaan stok.",
            ],
            assets["catalog"],
            "Antarmuka beranda dan katalog produk",
        ),
        (
            "Kesimpulan dan Saran",
            [
                "MotifNesia mendukung transformasi digital Toko Viera dari proses konvensional ke e-commerce.",
                "Sistem meningkatkan efisiensi stok, laporan, transaksi, dan layanan pelanggan.",
                "Katalog informatif, tracking, live chat, review, dan membership memperkuat pengalaman pelanggan.",
                "Pengembangan lanjutan disarankan pada perbandingan produk, konten edukasi batik, keamanan, performa, dan pemanfaatan data membership.",
            ],
            assets["supplier"],
            "Portal supplier sebagai bagian dari pengembangan operasional",
        ),
    ]

    for idx, (title, bullets, image, caption) in enumerate(slides, start=2):
        s = blank_slide(prs, idx)
        add_title(s, title, "MotifNesia")
        if image:
            add_bullets(s, bullets, Inches(0.75), Inches(1.65), Inches(5.15), Inches(4.95), 15.5)
            add_image(s, image, Inches(6.25), Inches(1.45), Inches(6.25), Inches(4.95), caption)
        else:
            add_rect(s, Inches(7.55), Inches(1.62), Inches(4.55), Inches(4.35), PAPER, RGBColor(235, 225, 214))
            add_bullets(s, bullets, Inches(0.9), Inches(1.75), Inches(6.25), Inches(4.8), 18)
            key = s.shapes.add_textbox(Inches(8.0), Inches(2.2), Inches(3.7), Inches(2.8))
            tf = key.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "Fokus"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = TERRACOTTA
            for line in ["Digitalisasi", "Efisiensi", "Transparansi", "Pengalaman Pelanggan"]:
                q = tf.add_paragraph()
                q.text = line
                q.font.size = Pt(24)
                q.font.bold = True
                q.font.color.rgb = INDIGO if line in ["Digitalisasi", "Transparansi"] else GREEN
                q.space_before = Pt(8)

    prs.save(PPTX)
    return PPTX


if __name__ == "__main__":
    path = build_deck()
    print(path)
