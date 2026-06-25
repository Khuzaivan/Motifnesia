from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PDF = Path(r"D:\Telyu\Semester 4\Proyek SI\coba.pdf")
OUT_DIR = ROOT / "generated"
ASSET_DIR = OUT_DIR / "story_assets"
PPTX = OUT_DIR / "MotifNesia_Presentasi_Storyline.pptx"

W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(251, 249, 244)
INK = RGBColor(38, 34, 31)
MUTED = RGBColor(105, 94, 86)
TERRACOTTA = RGBColor(164, 75, 48)
NAVY = RGBColor(31, 58, 95)
GREEN = RGBColor(70, 117, 91)
GOLD = RGBColor(204, 153, 78)
PAPER = RGBColor(255, 252, 246)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(226, 216, 204)
RED = RGBColor(177, 61, 54)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def render_clip(doc, page_no, bbox, name, zoom=3):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    page = doc[page_no - 1]
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), clip=fitz.Rect(bbox), alpha=False)
    path = ASSET_DIR / f"{name}.png"
    pix.save(path)
    return path


def render_assets():
    OUT_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(exist_ok=True)
    doc = fitz.open(PDF)
    return {
        "bpmn_as_is": render_clip(doc, 20, (105, 105, 520, 305), "bpmn_as_is_whatsapp"),
        "bpmn_to_be": render_clip(doc, 31, (105, 510, 555, 720), "bpmn_to_be_payment"),
        "use_case": render_clip(doc, 57, (105, 135, 545, 405), "use_case"),
        "architecture": render_clip(doc, 53, (160, 105, 495, 360), "architecture"),
        "homepage": render_clip(doc, 126, (108, 178, 515, 397), "homepage_catalog_1"),
        "catalog": render_clip(doc, 126, (108, 405, 515, 625), "homepage_catalog_2"),
        "checkout": render_clip(doc, 129, (164, 105, 460, 430), "checkout"),
        "admin_order": render_clip(doc, 138, (108, 105, 515, 340), "admin_order"),
        "sales_report": render_clip(doc, 139, (108, 105, 515, 575), "sales_report"),
    }


def add_bg(slide, n):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRACOTTA
    bar.line.color.rgb = TERRACOTTA
    foot = slide.shapes.add_textbox(Inches(12.15), Inches(7.08), Inches(0.65), Inches(0.22))
    set_run(foot.text_frame, f"{n:02d}", 8, MUTED, True, PP_ALIGN.RIGHT)


def set_run(tf, text, size, color=INK, bold=False, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(0)
    if align:
        p.alignment = align


def title(slide, text, kicker="MotifNesia"):
    k = slide.shapes.add_textbox(Inches(0.65), Inches(0.33), Inches(5.0), Inches(0.25))
    set_run(k.text_frame, kicker.upper(), 8.5, TERRACOTTA, True)
    t = slide.shapes.add_textbox(Inches(0.65), Inches(0.61), Inches(9.5), Inches(0.58))
    set_run(t.text_frame, text, 25, INK, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.23), Inches(0.95), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.color.rgb = GOLD


def bullets(slide, items, x, y, w, h, size=16, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def image(slide, path, x, y, w, h, caption=None):
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Inches(0.04), y - Inches(0.04), w + Inches(0.08), h + Inches(0.08))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if caption:
        c = slide.shapes.add_textbox(x, y + h + Inches(0.09), w, Inches(0.28))
        set_run(c.text_frame, caption, 8.5, MUTED)


def tag(slide, text, x, y, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.6), Inches(0.52))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(tf, text, 12, WHITE, True, PP_ALIGN.CENTER)
    return shp


def flow(slide, labels, x, y, box_w, box_h, gap, colors=None, size=11):
    colors = colors or [NAVY] * len(labels)
    shapes = []
    for i, label in enumerate(labels):
        bx = x + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, box_w, box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = colors[i % len(colors)]
        shp.line.color.rgb = colors[i % len(colors)]
        shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_run(shp.text_frame, label, size, WHITE, True, PP_ALIGN.CENTER)
        shapes.append(shp)
        if i < len(labels) - 1:
            arr = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                bx + box_w + Inches(0.04),
                y + box_h / 2,
                bx + box_w + gap - Inches(0.04),
                y + box_h / 2,
            )
            arr.line.color.rgb = MUTED
            arr.line.width = Pt(1.4)
            arr.line.end_arrowhead = True
    return shapes


def mini_card(slide, title_text, body, x, y, w, h, color=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = PAPER
    shp.line.color.rgb = LINE
    head = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.15), w - Inches(0.36), Inches(0.3))
    set_run(head.text_frame, title_text, 12.5, color, True)
    bullets(slide, body, x + Inches(0.18), y + Inches(0.55), w - Inches(0.36), h - Inches(0.65), 11.5, INK)


def new_slide(prs, n, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, n)
    title(slide, title_text)
    return slide


def build():
    a = render_assets()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 1)
    tag(s, "E-commerce Batik", Inches(0.72), Inches(0.7), TERRACOTTA)
    box = s.shapes.add_textbox(Inches(0.72), Inches(1.35), Inches(6.5), Inches(1.45))
    set_run(box.text_frame, "MotifNesia", 46, INK, True)
    sub = s.shapes.add_textbox(Inches(0.75), Inches(2.55), Inches(6.0), Inches(0.8))
    set_run(sub.text_frame, "Website penjualan online untuk membantu digitalisasi UMKM Batik Viera Garut", 20, MUTED)
    bullets(
        s,
        ["Fokus: pemasaran, transaksi, stok, laporan, dan layanan pelanggan", "Durasi presentasi: 10-15 menit"],
        Inches(0.78),
        Inches(4.45),
        Inches(5.8),
        Inches(1.0),
        15,
    )
    image(s, a["homepage"], Inches(7.2), Inches(0.8), Inches(5.25), Inches(4.85), "Homepage dan katalog sebagai wajah utama sistem")

    # 2
    s = new_slide(prs, 2, "Permasalahan UMKM Batik Viera")
    mini_card(s, "Pemasaran", ["Jangkauan masih bertumpu pada toko fisik", "Promosi belum terpusat secara digital"], Inches(0.75), Inches(1.65), Inches(3.65), Inches(1.65), TERRACOTTA)
    mini_card(s, "Transaksi", ["Pemesanan online masih melalui WhatsApp", "Bukti bayar diverifikasi manual"], Inches(4.85), Inches(1.65), Inches(3.65), Inches(1.65), NAVY)
    mini_card(s, "Operasional", ["Stok dan laporan perlu rekap manual", "Retur dan pengadaan belum terdokumentasi rapi"], Inches(8.95), Inches(1.65), Inches(3.65), Inches(1.65), GREEN)
    bullets(s, ["Dampak utama: informasi lambat, risiko salah stok, dan pengalaman belanja kurang transparan."], Inches(1.0), Inches(4.25), Inches(11.0), Inches(0.8), 21, RED)

    # 3
    s = new_slide(prs, 3, "Analisis Sistem Berjalan: As-Is")
    flow(
        s,
        ["Pelanggan\nbertanya", "Petugas\ncek stok", "Order via\nkasir/WA", "Bayar & kirim\nbukti", "Verifikasi\nmanual", "Barang\ndikirim"],
        Inches(0.7),
        Inches(1.65),
        Inches(1.65),
        Inches(0.9),
        Inches(0.28),
        [TERRACOTTA, NAVY, NAVY, GREEN, TERRACOTTA, GREEN],
    )
    bullets(s, ["BPMN As-Is menunjukkan proses masih bergantung pada petugas dan komunikasi manual.", "Pelanggan tidak mendapat status pesanan secara mandiri.", "Stok perlu dicocokkan kembali dengan catatan gudang."], Inches(0.85), Inches(3.05), Inches(5.4), Inches(2.8), 15)
    image(s, a["bpmn_as_is"], Inches(6.65), Inches(2.95), Inches(5.4), Inches(2.4), "BPMN As-Is: pemesanan non-tatap muka melalui WhatsApp")

    # 4
    s = new_slide(prs, 4, "Kesenjangan yang Harus Diselesaikan")
    bullets(
        s,
        ["Informasi stok belum real-time", "Checkout dan pembayaran belum terintegrasi", "Laporan penjualan belum langsung tersedia", "Retur belum berbasis bukti digital", "Promosi dan loyalitas pelanggan belum terstruktur"],
        Inches(0.9),
        Inches(1.55),
        Inches(5.4),
        Inches(4.7),
        20,
    )
    flow(s, ["Manual", "Lambat", "Tidak transparan", "Sulit diskalakan"], Inches(7.0), Inches(2.2), Inches(1.25), Inches(0.75), Inches(0.25), [RED, TERRACOTTA, NAVY, GREEN], 10.5)
    bullets(s, ["Arah solusi: satu platform yang menyatukan katalog, transaksi, stok, laporan, dan layanan pelanggan."], Inches(7.05), Inches(3.55), Inches(4.8), Inches(1.2), 17, NAVY)

    # 5
    s = new_slide(prs, 5, "Solusi To-Be: Website MotifNesia")
    flow(
        s,
        ["Katalog\nonline", "Keranjang &\ncheckout", "Upload bukti\nbayar", "Admin\nverifikasi", "Stok & status\nterbarui", "Laporan\notomatis"],
        Inches(0.7),
        Inches(1.58),
        Inches(1.65),
        Inches(0.9),
        Inches(0.28),
        [GREEN, NAVY, TERRACOTTA, TERRACOTTA, GREEN, NAVY],
    )
    bullets(s, ["Proses To-Be memindahkan aktivitas utama ke platform web.", "Pesanan yang melewati batas pembayaran dapat dibatalkan dan stok dikembalikan.", "Data transaksi menjadi dasar laporan, membership, dan evaluasi layanan."], Inches(0.85), Inches(3.05), Inches(5.35), Inches(2.6), 15)
    image(s, a["bpmn_to_be"], Inches(6.65), Inches(2.95), Inches(5.55), Inches(2.45), "BPMN To-Be: verifikasi pembayaran dan pesanan expired")

    # 6
    s = new_slide(prs, 6, "Perancangan Sistem: Aktor dan Use Case")
    bullets(s, ["Guest: melihat beranda, katalog, detail produk, dan registrasi.", "Customer: membeli produk, checkout, tracking, review, retur, live chat, dan membership.", "Admin: mengelola produk, pesanan, pembayaran, laporan, stok, retur, promo, dan konten.", "Supplier: memperbarui status pengadaan stok."], Inches(0.75), Inches(1.55), Inches(4.65), Inches(4.9), 14.5)
    image(s, a["use_case"], Inches(5.85), Inches(1.38), Inches(6.3), Inches(4.05), "Use Case Diagram sistem MotifNesia")

    # 7
    s = new_slide(prs, 7, "Arsitektur Sistem")
    image(s, a["architecture"], Inches(0.9), Inches(1.45), Inches(5.3), Inches(4.0), "Arsitektur: client, API, backend, database, dan layanan eksternal")
    bullets(s, ["Frontend web menjadi antarmuka customer, admin, dan supplier.", "API menghubungkan tampilan dengan layanan backend.", "Backend memproses login, produk, cart, checkout, pembayaran, pesanan, notifikasi, dan laporan.", "Database menyimpan data pengguna, produk, transaksi, stok, membership, dan supplier."], Inches(6.65), Inches(1.5), Inches(5.6), Inches(4.2), 16)

    # 8
    s = new_slide(prs, 8, "Fitur Utama untuk Customer")
    mini_card(s, "Belanja", ["Homepage", "Katalog dan filter", "Detail produk dan filosofi motif"], Inches(0.75), Inches(1.55), Inches(3.55), Inches(1.95), GREEN)
    mini_card(s, "Transaksi", ["Keranjang", "Checkout", "Konfirmasi pembayaran", "Tracking pesanan"], Inches(4.75), Inches(1.55), Inches(3.55), Inches(1.95), NAVY)
    mini_card(s, "Layanan", ["Retur digital", "Live chat", "Review produk", "Membership dan voucher"], Inches(8.75), Inches(1.55), Inches(3.55), Inches(1.95), TERRACOTTA)
    bullets(s, ["Inti pengalaman pelanggan: memilih produk lebih mudah, membeli lebih praktis, dan memantau pesanan dengan transparan."], Inches(1.0), Inches(4.55), Inches(11.0), Inches(0.8), 20, INK)

    # 9
    s = new_slide(prs, 9, "Fitur Utama untuk Admin dan Supplier")
    mini_card(s, "Admin Penjualan", ["Kelola produk", "Verifikasi pembayaran", "Status pesanan", "Laporan penjualan"], Inches(0.75), Inches(1.55), Inches(3.65), Inches(2.1), NAVY)
    mini_card(s, "Admin Operasional", ["Stok gudang", "Stock opname", "Pengadaan stok", "Manajemen retur"], Inches(4.85), Inches(1.55), Inches(3.65), Inches(2.1), GREEN)
    mini_card(s, "Supplier", ["Daftar permintaan stok", "Detail pengadaan", "Update status barang"], Inches(8.95), Inches(1.55), Inches(3.65), Inches(2.1), TERRACOTTA)
    bullets(s, ["Tujuannya bukan hanya menjual online, tetapi membuat operasional toko lebih terukur."], Inches(1.0), Inches(4.65), Inches(10.7), Inches(0.75), 20, INK)

    # 10
    s = new_slide(prs, 10, "Implementasi: Homepage dan Katalog")
    image(s, a["homepage"], Inches(0.75), Inches(1.45), Inches(5.55), Inches(3.05), "Homepage: akses awal, kategori, dan produk unggulan")
    image(s, a["catalog"], Inches(6.85), Inches(1.45), Inches(5.55), Inches(3.05), "Katalog: kartu produk, harga, promo, dan tombol detail")
    bullets(s, ["Slide ini menunjukkan bagian yang pertama kali memengaruhi keputusan pelanggan: tampilan produk dan kemudahan memilih."], Inches(1.0), Inches(5.25), Inches(11.0), Inches(0.75), 17)

    # 11
    s = new_slide(prs, 11, "Implementasi: Checkout dan Pembayaran")
    image(s, a["checkout"], Inches(0.9), Inches(1.35), Inches(4.15), Inches(4.65), "Checkout")
    bullets(s, ["Pelanggan memilih alamat, pengiriman, dan metode pembayaran.", "Sistem menampilkan ringkasan produk dan total pembayaran.", "Bukti pembayaran masuk ke proses verifikasi admin.", "Alur ini menjadi pengganti pemesanan WhatsApp yang sebelumnya manual."], Inches(5.65), Inches(1.6), Inches(5.9), Inches(3.7), 18)

    # 12
    s = new_slide(prs, 12, "Implementasi: Dashboard Admin dan Laporan")
    image(s, a["admin_order"], Inches(0.7), Inches(1.45), Inches(5.65), Inches(3.1), "Dashboard status pesanan")
    image(s, a["sales_report"], Inches(6.75), Inches(1.45), Inches(5.65), Inches(3.35), "Laporan penjualan")
    bullets(s, ["Admin dapat memantau pesanan, bukti pembayaran, status pengiriman, omzet, produk terjual, dan riwayat transaksi."], Inches(1.0), Inches(5.45), Inches(11.0), Inches(0.75), 17)

    # 13
    s = new_slide(prs, 13, "Alur Penggunaan Sistem")
    flow(
        s,
        ["Customer\npilih produk", "Checkout &\nbayar", "Admin\nverifikasi", "Pesanan\ndiproses", "Stok & laporan\nterupdate", "Customer\nreview/retur"],
        Inches(0.7),
        Inches(1.55),
        Inches(1.65),
        Inches(0.95),
        Inches(0.28),
        [GREEN, NAVY, TERRACOTTA, NAVY, GREEN, TERRACOTTA],
    )
    bullets(s, ["Alur ini memperlihatkan hubungan langsung antara aktivitas pelanggan dan pekerjaan admin.", "Setiap transaksi berdampak pada stok, status pesanan, laporan, poin membership, dan evaluasi layanan.", "Supplier masuk ketika stok menipis dan admin membuat pengadaan."], Inches(1.0), Inches(3.35), Inches(11.0), Inches(2.0), 18)

    # 14
    s = new_slide(prs, 14, "Hasil Pengujian Sistem")
    bullets(s, ["Skenario use case telah disusun untuk fitur utama: login, katalog, checkout, transaksi, pesanan, stok, laporan, retur, chat, membership, dan supplier.", "Validasi yang dirancang mencakup input benar/salah, stok tersedia/habis, bukti pembayaran valid/tidak valid, kode promo, dan export laporan.", "Catatan laporan: bagian 4.2 Pengujian masih belum memuat tabel hasil uji formal, sehingga dokumentasi pengujian perlu dilengkapi sebelum sidang final."], Inches(0.85), Inches(1.55), Inches(6.35), Inches(4.6), 17)
    mini_card(s, "Status", ["Fitur utama sudah diimplementasikan pada antarmuka", "Skenario pengujian sudah tergambar pada use case", "Tabel hasil uji formal belum lengkap di PDF"], Inches(8.0), Inches(1.75), Inches(3.75), Inches(2.4), TERRACOTTA)

    # 15
    s = new_slide(prs, 15, "Kesimpulan dan Manfaat")
    mini_card(s, "Bagi Toko Viera", ["Pemasaran lebih luas", "Stok dan laporan lebih terkontrol", "Pengadaan supplier lebih tercatat"], Inches(0.75), Inches(1.55), Inches(3.7), Inches(2.35), TERRACOTTA)
    mini_card(s, "Bagi Pelanggan", ["Informasi produk lebih jelas", "Transaksi lebih praktis", "Pesanan lebih transparan"], Inches(4.85), Inches(1.55), Inches(3.7), Inches(2.35), GREEN)
    mini_card(s, "Pengembangan Lanjut", ["Perbandingan produk", "Konten edukasi batik", "Peningkatan keamanan dan performa", "Analisis data membership"], Inches(8.95), Inches(1.55), Inches(3.7), Inches(2.35), NAVY)
    bullets(s, ["MotifNesia bukan hanya katalog online, tetapi fondasi transformasi digital untuk proses penjualan dan operasional UMKM batik."], Inches(1.0), Inches(4.8), Inches(11.0), Inches(0.9), 21, INK)

    prs.save(PPTX)
    print(PPTX)


if __name__ == "__main__":
    build()
