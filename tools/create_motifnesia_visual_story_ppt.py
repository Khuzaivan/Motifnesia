from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PDF = Path(r"D:\Telyu\Semester 4\Proyek SI\coba.pdf")
OUT_DIR = ROOT / "generated"
ASSET_DIR = OUT_DIR / "visual_story_assets"
PPTX = OUT_DIR / "MotifNesia_Presentasi_Visual_Storytelling.pptx"

W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(250, 248, 243)
INK = RGBColor(34, 31, 29)
MUTED = RGBColor(103, 91, 84)
BROWN = RGBColor(135, 74, 50)
RUST = RGBColor(174, 82, 49)
NAVY = RGBColor(34, 60, 95)
GREEN = RGBColor(62, 113, 91)
GOLD = RGBColor(203, 151, 72)
RED = RGBColor(174, 61, 55)
WHITE = RGBColor(255, 255, 255)
PAPER = RGBColor(255, 252, 246)
LINE = RGBColor(226, 216, 204)


def render_clip(doc, page_no, bbox, name, zoom=3.3):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    page = doc[page_no - 1]
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), clip=fitz.Rect(bbox), alpha=False)
    path = ASSET_DIR / f"{name}.png"
    pix.save(path)
    return path


def render_assets():
    OUT_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(exist_ok=True)
    doc = fitz.open(PDF)
    return {
        "as_is": render_clip(doc, 20, (105, 105, 520, 305), "bpmn_as_is"),
        "to_be": render_clip(doc, 31, (105, 510, 555, 720), "bpmn_to_be_payment"),
        "architecture": render_clip(doc, 53, (160, 105, 495, 360), "architecture"),
        "use_case": render_clip(doc, 57, (105, 135, 545, 405), "use_case"),
        "home_top": render_clip(doc, 126, (108, 178, 515, 397), "home_top"),
        "catalog": render_clip(doc, 126, (108, 405, 515, 625), "catalog"),
        "detail": render_clip(doc, 127, (108, 112, 515, 625), "detail_product"),
        "cart": render_clip(doc, 128, (108, 112, 515, 440), "cart"),
        "checkout": render_clip(doc, 129, (164, 105, 460, 430), "checkout"),
        "payment": render_clip(doc, 130, (108, 105, 515, 495), "payment_confirm"),
        "admin_order": render_clip(doc, 138, (108, 105, 515, 340), "admin_order"),
        "sales_report": render_clip(doc, 139, (108, 105, 515, 575), "sales_report"),
        "membership": render_clip(doc, 132, (108, 105, 515, 610), "membership"),
        "retur": render_clip(doc, 134, (108, 105, 515, 575), "retur"),
        "chat": render_clip(doc, 135, (108, 105, 515, 575), "chat"),
        "stock": render_clip(doc, 146, (108, 105, 515, 575), "stock"),
    }


def text_frame(shape, text, size=18, color=INK, bold=False, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(0)
    if align:
        p.alignment = align
    return tf


def bg(slide, n):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BROWN
    bar.line.color.rgb = BROWN
    num = slide.shapes.add_textbox(Inches(12.33), Inches(7.06), Inches(0.55), Inches(0.24))
    text_frame(num, f"{n:02d}", 8, MUTED, True, PP_ALIGN.RIGHT)


def headline(slide, msg, kicker=None):
    if kicker:
        k = slide.shapes.add_textbox(Inches(0.65), Inches(0.34), Inches(4.5), Inches(0.24))
        text_frame(k, kicker.upper(), 8.5, RUST, True)
    h = slide.shapes.add_textbox(Inches(0.65), Inches(0.68), Inches(9.6), Inches(0.8))
    text_frame(h, msg, 26, INK, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(1.45), Inches(0.95), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.color.rgb = GOLD


def bullets(slide, items, x, y, w, h, size=15, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_image(slide, path, x, y, w, h):
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Inches(0.05), y - Inches(0.05), w + Inches(0.1), h + Inches(0.1))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def caption(slide, txt, x, y, w):
    c = slide.shapes.add_textbox(x, y, w, Inches(0.26))
    text_frame(c, txt, 8.5, MUTED)


def callout(slide, txt, x, y, w, color=NAVY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.62))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame(box, txt, 12.5, WHITE, True, PP_ALIGN.CENTER)
    return box


def stat(slide, number, label, x, y, color):
    n = slide.shapes.add_textbox(x, y, Inches(2.1), Inches(0.7))
    text_frame(n, number, 32, color, True, PP_ALIGN.CENTER)
    l = slide.shapes.add_textbox(x, y + Inches(0.68), Inches(2.1), Inches(0.44))
    text_frame(l, label, 10.5, MUTED, False, PP_ALIGN.CENTER)


def card(slide, title, items, x, y, w, h, color=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = PAPER
    shp.line.color.rgb = LINE
    t = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.3))
    text_frame(t, title, 12.5, color, True)
    bullets(slide, items, x + Inches(0.18), y + Inches(0.58), w - Inches(0.36), h - Inches(0.64), 11.4)


def flow(slide, labels, x, y, box_w, gap):
    colors = [GREEN, NAVY, RUST, NAVY, GREEN]
    for i, label in enumerate(labels):
        bx = x + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, box_w, Inches(0.82))
        shp.fill.solid()
        shp.fill.fore_color.rgb = colors[i % len(colors)]
        shp.line.color.rgb = colors[i % len(colors)]
        shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame(shp, label, 10.7, WHITE, True, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                bx + box_w + Inches(0.03),
                y + Inches(0.41),
                bx + box_w + gap - Inches(0.03),
                y + Inches(0.41),
            )
            conn.line.color.rgb = MUTED
            conn.line.width = Pt(1.2)
            conn.line.end_arrowhead = True


def slide(prs, n, msg, kicker="Cerita Proyek"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, n)
    headline(s, msg, kicker)
    return s


def build():
    a = render_assets()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, 1)
    callout(s, "UMKM Batik Viera Garut", Inches(0.72), Inches(0.78), Inches(2.9), RUST)
    h = s.shapes.add_textbox(Inches(0.72), Inches(1.55), Inches(6.05), Inches(1.55))
    text_frame(h, "MotifNesia mengubah toko batik menjadi pengalaman belanja digital", 33, INK, True)
    bullets(s, ["Website e-commerce untuk katalog, transaksi, stok, laporan, dan layanan pelanggan"], Inches(0.78), Inches(4.65), Inches(5.8), Inches(0.75), 17)
    add_image(s, a["home_top"], Inches(7.15), Inches(0.82), Inches(5.25), Inches(4.85))
    caption(s, "Tampilan awal website MotifNesia", Inches(7.15), Inches(5.82), Inches(5.2))

    # 2
    s = slide(prs, 2, "Masalahnya bukan cuma belum online, tapi proses bisnisnya belum terhubung")
    stat(s, "Manual", "stok dan laporan", Inches(0.95), Inches(1.85), RED)
    stat(s, "WA", "pemesanan jarak jauh", Inches(3.55), Inches(1.85), RUST)
    stat(s, "Tidak real-time", "status dan ketersediaan", Inches(6.15), Inches(1.85), NAVY)
    stat(s, "Terpisah", "retur, supplier, promosi", Inches(8.9), Inches(1.85), GREEN)
    bullets(s, ["Akibatnya pelanggan perlu banyak bertanya, admin perlu banyak mengecek, dan owner sulit membaca kondisi bisnis secara cepat."], Inches(1.0), Inches(4.25), Inches(11.0), Inches(0.9), 20)

    # 3
    s = slide(prs, 3, "Pada proses As-Is, pelanggan masih bergantung pada petugas")
    add_image(s, a["as_is"], Inches(0.75), Inches(1.65), Inches(6.55), Inches(3.1))
    bullets(s, ["Cek stok dilakukan manual", "Bukti pembayaran dikirim lewat WhatsApp", "Status pengiriman tidak bisa dipantau mandiri"], Inches(7.75), Inches(1.95), Inches(4.3), Inches(2.4), 18)
    callout(s, "Risiko utama: lambat, rawan selisih, dan kurang transparan", Inches(7.75), Inches(4.7), Inches(4.25), RED)
    caption(s, "BPMN As-Is: pemesanan non-tatap muka", Inches(0.75), Inches(4.88), Inches(6.4))

    # 4
    s = slide(prs, 4, "MotifNesia menyatukan alur belanja dari katalog sampai laporan")
    add_image(s, a["to_be"], Inches(6.9), Inches(1.65), Inches(5.35), Inches(2.7))
    flow(s, ["Lihat\nproduk", "Checkout", "Upload\nbukti", "Admin\nverifikasi", "Laporan\nterbentuk"], Inches(0.75), Inches(2.05), Inches(1.35), Inches(0.26))
    bullets(s, ["Stok tampil lebih akurat", "Pesanan expired dapat dibatalkan otomatis", "Transaksi menjadi data laporan"], Inches(0.95), Inches(3.65), Inches(5.25), Inches(2.0), 17)
    caption(s, "BPMN To-Be: verifikasi pembayaran dan pesanan expired", Inches(6.9), Inches(4.52), Inches(5.2))

    # 5
    s = slide(prs, 5, "Rancangan sistem dibuat sederhana: Customer, Admin, dan Supplier")
    add_image(s, a["use_case"], Inches(0.7), Inches(1.62), Inches(6.15), Inches(3.75))
    card(s, "Aktor utama", ["Customer membeli dan memantau pesanan", "Admin mengelola operasional", "Supplier memperbarui pengadaan stok"], Inches(7.35), Inches(1.75), Inches(4.35), Inches(2.25), NAVY)
    card(s, "Pesan slide", ["Use case memastikan setiap fitur punya pemilik proses yang jelas"], Inches(7.35), Inches(4.25), Inches(4.35), Inches(1.35), GREEN)

    # 6
    s = slide(prs, 6, "Arsitektur memisahkan tampilan, logika bisnis, dan data")
    add_image(s, a["architecture"], Inches(0.85), Inches(1.55), Inches(5.8), Inches(4.15))
    bullets(s, ["Frontend: web customer, admin, supplier", "API: jalur komunikasi sistem", "Backend: autentikasi, transaksi, stok, laporan", "Database: produk, user, pesanan, membership"], Inches(7.15), Inches(1.8), Inches(4.85), Inches(3.35), 17)
    callout(s, "Hasilnya: sistem lebih mudah dikembangkan dan diintegrasikan", Inches(7.15), Inches(5.15), Inches(4.65), GREEN)

    # 7
    s = slide(prs, 7, "Customer pertama kali diyakinkan lewat homepage dan katalog")
    add_image(s, a["home_top"], Inches(0.7), Inches(1.6), Inches(5.75), Inches(3.15))
    add_image(s, a["catalog"], Inches(6.9), Inches(1.6), Inches(5.75), Inches(3.15))
    bullets(s, ["Produk terlihat langsung", "Harga dan promo mudah dibandingkan", "Navigasi dibuat untuk memilih produk cepat"], Inches(1.0), Inches(5.3), Inches(10.5), Inches(0.85), 16)

    # 8
    s = slide(prs, 8, "Detail produk menjawab keraguan pelanggan sebelum membeli")
    add_image(s, a["detail"], Inches(0.8), Inches(1.55), Inches(6.0), Inches(4.55))
    bullets(s, ["Foto produk", "Ukuran dan stok", "Material dan proses pembuatan", "Deskripsi, filosofi motif, dan ulasan"], Inches(7.35), Inches(1.8), Inches(4.3), Inches(2.6), 18)
    callout(s, "Fokus UX: membangun trust pada kualitas batik", Inches(7.35), Inches(4.75), Inches(4.35), RUST)

    # 9
    s = slide(prs, 9, "Checkout menggantikan alur pemesanan manual lewat chat")
    add_image(s, a["checkout"], Inches(0.9), Inches(1.5), Inches(4.35), Inches(4.85))
    add_image(s, a["payment"], Inches(5.85), Inches(1.5), Inches(4.35), Inches(4.85))
    bullets(s, ["Pilih alamat dan pengiriman", "Pilih metode pembayaran", "Upload bukti pembayaran", "Admin memverifikasi dari dashboard"], Inches(10.65), Inches(1.75), Inches(1.85), Inches(3.4), 13.5)

    # 10
    s = slide(prs, 10, "Admin punya pusat kendali untuk pesanan dan laporan")
    add_image(s, a["admin_order"], Inches(0.75), Inches(1.55), Inches(5.65), Inches(3.3))
    add_image(s, a["sales_report"], Inches(6.85), Inches(1.55), Inches(5.65), Inches(3.65))
    bullets(s, ["Pesanan masuk terlihat terpusat", "Pembayaran dan pengiriman bisa dipantau", "Laporan omzet, produk terjual, dan riwayat transaksi tersedia"], Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.95), 16)

    # 11
    s = slide(prs, 11, "Fitur pendukung membuat pelanggan kembali dan layanan lebih rapi")
    add_image(s, a["membership"], Inches(0.75), Inches(1.58), Inches(3.6), Inches(3.9))
    add_image(s, a["retur"], Inches(4.85), Inches(1.58), Inches(3.6), Inches(3.9))
    add_image(s, a["chat"], Inches(8.95), Inches(1.58), Inches(3.6), Inches(3.9))
    bullets(s, ["Membership menjaga loyalitas", "Retur terdokumentasi digital", "Live chat mempercepat bantuan pelanggan"], Inches(1.0), Inches(5.85), Inches(11.0), Inches(0.7), 16)

    # 12
    s = slide(prs, 12, "Alur akhirnya sederhana: pelanggan belanja, admin mengelola, data menjadi keputusan")
    flow(s, ["Customer\npilih produk", "Checkout &\nbayar", "Admin\nverifikasi", "Stok/laporan\nterupdate", "Evaluasi\nlayanan"], Inches(1.1), Inches(2.1), Inches(1.75), Inches(0.5))
    bullets(s, ["Satu transaksi menggerakkan beberapa proses sekaligus: stok, status pesanan, laporan, poin membership, dan review layanan.", "Inilah perbedaan utama MotifNesia dibanding proses lama yang terpisah-pisah."], Inches(1.2), Inches(4.05), Inches(10.6), Inches(1.35), 19)

    # 13
    s = slide(prs, 13, "Pengujian diarahkan ke skenario bisnis yang paling berisiko")
    callout(s, "Transaksi", Inches(1.0), Inches(2.0), Inches(2.65), NAVY)
    callout(s, "Operasional", Inches(5.3), Inches(2.0), Inches(2.65), GREEN)
    callout(s, "Layanan", Inches(9.6), Inches(2.0), Inches(2.65), RUST)
    bullets(
        s,
        [
            "Validasi difokuskan pada checkout, pembayaran, stok, laporan, retur, live chat, membership, dan supplier.",
            "Skenario mencakup kondisi berhasil dan gagal, misalnya bukti bayar tidak valid atau stok tidak tersedia.",
            "Catatan: bagian hasil pengujian formal pada PDF masih perlu dilengkapi sebelum sidang final.",
        ],
        Inches(1.1),
        Inches(3.45),
        Inches(10.75),
        Inches(1.6),
        18,
        INK,
    )

    # 14
    s = slide(prs, 14, "Manfaat utamanya: toko lebih siap masuk pasar digital")
    stat(s, "Lebih luas", "jangkauan pemasaran", Inches(1.0), Inches(1.85), RUST)
    stat(s, "Lebih cepat", "transaksi dan layanan", Inches(3.85), Inches(1.85), NAVY)
    stat(s, "Lebih rapi", "stok, retur, supplier", Inches(6.7), Inches(1.85), GREEN)
    stat(s, "Lebih terukur", "laporan dan evaluasi", Inches(9.55), Inches(1.85), GOLD)
    bullets(s, ["MotifNesia bukan sekadar etalase online, tetapi sistem yang menghubungkan penjualan, operasional, dan pengambilan keputusan."], Inches(1.0), Inches(4.55), Inches(11.0), Inches(0.9), 21)

    prs.save(PPTX)
    print(PPTX)


if __name__ == "__main__":
    build()
